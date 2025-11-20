import numpy as np
import pandas as pd
import streamlit as st
import plotly.graph_objects as go
import io

# Optional exports
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False


# -------------------------------------------------------------------
# 1. Synthetic data setup
# -------------------------------------------------------------------

LOB_LIST = [
    "Technology & Infrastructure",
    "Real Estate & Facilities",
    "Finance & Treasury",
    "Human Resources",
    "Legal & Compliance",
]

np.random.seed(42)


def _base_headcount(lob: str) -> int:
    mapping = {
        "Technology & Infrastructure": 8000,
        "Real Estate & Facilities": 3000,
        "Finance & Treasury": 2000,
        "Human Resources": 1500,
        "Legal & Compliance": 1200,
    }
    return mapping[lob]


def _avg_compensation(lob: str) -> float:
    high = ["Technology & Infrastructure", "Legal & Compliance"]
    return 140_000 if lob in high else 110_000


@st.cache_data
def generate_historical_data(start_q: str = "2020Q1", end_q: str = "2024Q4") -> pd.DataFrame:
    """
    Create synthetic quarterly data:
    - Macro: GDP, unemployment, equity index
    - LOB-level: headcount, expense buckets, recoveries, capital, RWA
    """
    periods = pd.period_range(start_q, end_q, freq="Q")
    rows = []

    for p in periods:
        year, q = p.year, p.quarter

        # Simple synthetic macro drivers
        gdp_growth = 2 + 0.5 * np.sin((p.qyear) / 2) + np.random.normal(0, 0.3)
        unemployment = 4.5 + 0.3 * np.cos((p.qyear) / 3) + np.random.normal(0, 0.2)
        equity_index = 100 + (p.qyear - 2020) * 5 + np.random.normal(0, 3)

        for lob in LOB_LIST:
            base_hc = _base_headcount(lob)
            trend = (year - 2020) * 0.03
            headcount = int(base_hc * (1 + trend) * (1 + np.random.normal(0, 0.02)))

            avg_comp = _avg_compensation(lob)
            salary_expense = headcount * avg_comp / 4  # quarterly

            vendor_expense = salary_expense * 0.30 * (1 + np.random.normal(0, 0.10))
            real_estate_expense = base_hc * 5_000 * (1 + 0.02 * (year - 2020))
            other_expense = salary_expense * 0.10 * (1 + np.random.normal(0, 0.10))

            total_expense = (
                salary_expense
                + vendor_expense
                + real_estate_expense
                + other_expense
            )

            # Internal recoveries / allocated revenue
            allocated_revenue = total_expense * (1.05 + np.random.normal(0, 0.05))
            pretax_income = allocated_revenue - total_expense

            # Simple capital / RWA proxy
            rwa = total_expense / 1_000
            capital = rwa * 0.13 + np.random.normal(0, rwa * 0.01)
            capital_ratio = capital / rwa

            rows.append(
                {
                    "quarter": str(p),
                    "year": year,
                    "qtr": q,
                    "lob": lob,
                    "gdp_growth": gdp_growth,
                    "unemployment": unemployment,
                    "equity_index": equity_index,
                    "headcount": headcount,
                    "salary_expense": salary_expense,
                    "vendor_expense": vendor_expense,
                    "real_estate_expense": real_estate_expense,
                    "other_expense": other_expense,
                    "total_expense": total_expense,
                    "allocated_revenue": allocated_revenue,
                    "pretax_income": pretax_income,
                    "rwa": rwa,
                    "capital": capital,
                    "capital_ratio": capital_ratio,
                }
            )

    return pd.DataFrame(rows)


# -------------------------------------------------------------------
# 2. Scenario engine
# -------------------------------------------------------------------

def project_scenario(
    df_hist: pd.DataFrame,
    scenario_name: str,
    rev_shock_yr: float,
    exp_shock_yr: float,
    horizon_quarters: int,
    capital_floor_ratio: float,
) -> pd.DataFrame:
    """
    Apply a simple satellite-model style projection:
    - Annual shocks to revenue / expense -> quarterly growth rates
    - Losses reduce capital
    - Capital floored at specified capital_floor_ratio * RWA
    """
    df_hist = df_hist.copy()
    df_hist["quarter"] = pd.PeriodIndex(df_hist["quarter"], freq="Q")
    last_q = df_hist["quarter"].max()

    base = (
        df_hist[df_hist["quarter"] == last_q]
        .groupby("lob")
        .agg(
            headcount=("headcount", "sum"),
            total_expense=("total_expense", "sum"),
            allocated_revenue=("allocated_revenue", "sum"),
            capital=("capital", "sum"),
            rwa=("rwa", "sum"),
        )
        .reset_index()
    )
    base["capital_ratio"] = base["capital"] / base["rwa"]

    # Annual shocks -> quarterly growth
    rev_q = (1 + rev_shock_yr) ** (1 / 4) - 1
    exp_q = (1 + exp_shock_yr) ** (1 / 4) - 1

    rows = []
    curr_period = last_q
    current = base.copy()

    for _ in range(1, horizon_quarters + 1):
        curr_period += 1

        current["allocated_revenue"] *= (1 + rev_q)
        current["total_expense"] *= (1 + exp_q)
        current["headcount"] *= (1 + exp_q * 0.3)

        pretax = current["allocated_revenue"] - current["total_expense"]
        current["rwa"] = current["total_expense"] / 1_000

        # Losses eat capital
        loss = -pretax.clip(upper=0)
        current["capital"] = current["capital"] - loss * 0.7

        # Apply capital floor
        floor_capital = current["rwa"] * capital_floor_ratio
        current["capital"] = np.maximum(current["capital"], floor_capital)
        current["capital_ratio"] = current["capital"] / current["rwa"]

        for _, r in current.iterrows():
            rows.append(
                {
                    "quarter": str(curr_period),
                    "lob": r["lob"],
                    "headcount": r["headcount"],
                    "total_expense": r["total_expense"],
                    "allocated_revenue": r["allocated_revenue"],
                    "pretax_income": r["allocated_revenue"] - r["total_expense"],
                    "rwa": r["rwa"],
                    "capital": r["capital"],
                    "capital_ratio": r["capital_ratio"],
                    "scenario": scenario_name,
                }
            )

    return pd.DataFrame(rows)


def build_all_scenarios(
    df_hist: pd.DataFrame,
    scenario_params: dict,
    horizon_quarters: int,
    capital_floor_ratio: float,
) -> pd.DataFrame:
    """Run all scenarios and stack results."""
    dfs = []
    for name, params in scenario_params.items():
        dfs.append(
            project_scenario(
                df_hist,
                scenario_name=name,
                rev_shock_yr=params["rev_shock_yr"],
                exp_shock_yr=params["exp_shock_yr"],
                horizon_quarters=horizon_quarters,
                capital_floor_ratio=capital_floor_ratio,
            )
        )
    return pd.concat(dfs, ignore_index=True)


# -------------------------------------------------------------------
# 3. Metrics & visuals
# -------------------------------------------------------------------

def summarize_scenario(
    df_proj: pd.DataFrame,
    scenario: str,
    cap_threshold: float,
    eff_threshold: float,
):
    """
    Aggregate projections to firm level for a given scenario and compute:
    - capital ratio
    - efficiency ratio
    - cumulative losses
    - # of risk appetite breaches (capital OR efficiency)
    """
    df_scen = df_proj[df_proj["scenario"] == scenario].copy()

    agg = (
        df_scen.groupby("quarter")
        .agg(
            total_expense=("total_expense", "sum"),
            allocated_revenue=("allocated_revenue", "sum"),
            pretax_income=("pretax_income", "sum"),
            capital=("capital", "sum"),
            rwa=("rwa", "sum"),
        )
        .reset_index()
    )

    agg["capital_ratio"] = agg["capital"] / agg["rwa"]
    agg["efficiency_ratio"] = agg["total_expense"] / agg["allocated_revenue"]

    # Count a breach when capital ratio is below its minimum
    # OR efficiency ratio is above its maximum
    capital_breaches = agg["capital_ratio"] < cap_threshold
    efficiency_breaches = agg["efficiency_ratio"] > eff_threshold
    breaches = int((capital_breaches | efficiency_breaches).sum())

    cumulative_loss = (-agg["pretax_income"].clip(upper=0)).sum()

    metrics = {
        "cumulative_loss": float(cumulative_loss),
        "min_capital_ratio": float(agg["capital_ratio"].min()),
        "breaches": int(breaches),
    }

    return agg, metrics


def plot_capital_ratio(agg: pd.DataFrame, cap_threshold: float, scenario: str):
    fig = go.Figure()

    # Main capital ratio path
    fig.add_trace(
        go.Scatter(
            x=agg["quarter"],
            y=agg["capital_ratio"],
            mode="lines+markers",
            name="Capital Ratio",
        )
    )

    # Red markers on capital-ratio breach quarters
    breach_points = agg[agg["capital_ratio"] < cap_threshold]
    if not breach_points.empty:
        fig.add_trace(
            go.Scatter(
                x=breach_points["quarter"],
                y=breach_points["capital_ratio"],
                mode="markers",
                marker=dict(color="red", size=10),
                name="Capital Breaches",
            )
        )

    fig.add_hline(
        y=cap_threshold,
        line_dash="dash",
        annotation_text="Risk Appetite Threshold",
        annotation_position="top left",
    )
    fig.update_layout(
        title=f"Capital Ratio Path – {scenario}",
        xaxis_title="Quarter",
        yaxis_title="Capital Ratio",
    )
    return fig


def plot_capital_ratio_multi(agg_dict: dict, cap_threshold: float):
    fig = go.Figure()
    for scenario, agg in agg_dict.items():
        fig.add_trace(
            go.Scatter(
                x=agg["quarter"],
                y=agg["capital_ratio"],
                mode="lines+markers",
                name=scenario,
            )
        )
    fig.add_hline(
        y=cap_threshold,
        line_dash="dash",
        annotation_text="Risk Appetite Threshold",
        annotation_position="top left",
    )
    fig.update_layout(
        title="Capital Ratio Paths – Scenario Comparison",
        xaxis_title="Quarter",
        yaxis_title="Capital Ratio",
    )
    return fig


def plot_pretax_income_by_lob(df_proj: pd.DataFrame, scenario: str):
    df_scen = df_proj[df_proj["scenario"] == scenario]
    lob_pnl = (
        df_scen.groupby(["quarter", "lob"])["pretax_income"]
        .sum()
        .reset_index()
    )

    fig = go.Figure()
    for lob in LOB_LIST:
        sub = lob_pnl[lob_pnl["lob"] == lob]
        fig.add_trace(
            go.Bar(
                x=sub["quarter"],
                y=sub["pretax_income"],
                name=lob,
            )
        )

    fig.update_layout(
        barmode="relative",
        title=f"Pretax Income by Function – {scenario}",
        xaxis_title="Quarter",
        yaxis_title="Pretax Income",
    )
    return fig


# -------------------------------------------------------------------
# 4. PowerPoint & PDF exports
# -------------------------------------------------------------------

def build_ppt_report(agg_dict: dict, metrics_dict: dict) -> bytes:
    if not PPTX_AVAILABLE:
        raise RuntimeError("python-pptx is not installed")

    prs = Presentation()

    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Stress Testing & Scenario Analysis Summary"
    subtitle = slide.placeholders[1]
    subtitle.text = "Auto-generated report from dashboard"

    # Metrics slide
    layout = prs.slide_layouts[5]
    slide2 = prs.slides.add_slide(layout)
    title = slide2.shapes.title
    title.text = "Scenario Metrics Overview"

    rows = len(metrics_dict) + 1
    cols = 4
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(2.5)
    table = slide2.shapes.add_table(rows, cols, left, top, width, height).table

    headers = ["Scenario", "Cumulative Loss", "Min Capital Ratio", "# Breaches"]
    for j, h in enumerate(headers):
        table.cell(0, j).text = h

    for i, (scenario, m) in enumerate(metrics_dict.items(), start=1):
        table.cell(i, 0).text = scenario
        table.cell(i, 1).text = f"${m['cumulative_loss']:,.0f}"
        table.cell(i, 2).text = f"{m['min_capital_ratio']:.2%}"
        table.cell(i, 3).text = str(m["breaches"])

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


def build_pdf_report(
    metrics_dict: dict,
    cap_threshold: float,
    eff_threshold: float,
) -> bytes:
    if not PDF_AVAILABLE:
        raise RuntimeError("reportlab is not installed")

    buf = io.BytesIO()
    c = canvas.Canvas(buf, pagesize=letter)
    width, height = letter
    y = height - 50

    c.setFont("Helvetica-Bold", 16)
    c.drawString(50, y, "Stress Testing & Scenario Analysis Summary")

    y -= 30
    c.setFont("Helvetica", 10)
    c.drawString(50, y, f"Capital ratio threshold: {cap_threshold:.2%}")
    y -= 15
    c.drawString(50, y, f"Efficiency ratio threshold: {eff_threshold:.2%}")

    y -= 30
    c.setFont("Helvetica-Bold", 12)
    c.drawString(50, y, "Scenario Metrics:")
    y -= 20
    c.setFont("Helvetica", 10)

    for scenario, m in metrics_dict.items():
        if y < 80:
            c.showPage()
            y = height - 50
            c.setFont("Helvetica", 10)

        line = (
            f"{scenario}: "
            f"Cumulative loss = ${m['cumulative_loss']:,.0f}, "
            f"Min capital ratio = {m['min_capital_ratio']:.2%}, "
            f"Breaches = {m['breaches']}"
        )
        c.drawString(50, y, line)
        y -= 15

    c.showPage()
    c.save()
    buf.seek(0)
    return buf.getvalue()


# -------------------------------------------------------------------
# 5. Documentation tab
# -------------------------------------------------------------------

def make_documentation_tab():
    st.markdown(
        """
        ### Model Overview

        This sandbox models multi-scenario financial stress testing and forecasting
        over a multi-quarter horizon.

        **Key components:**
        - Synthetic historical data for multiple support functions (headcount, expenses, recoveries, capital, RWA).
        - Scenario engine projecting revenue, expenses, and capital under different stress assumptions.
        - Risk metrics including capital ratio, efficiency ratio, cumulative losses, and breach counts.

        ### Methodology (High-Level)

        1. **Historical Baseline**
           - Quarterly synthetic data are generated for several years.
           - For each function, we estimate headcount and derive compensation and non-compensation expenses.
           - Recoveries/allocated revenue are modeled as a markup over expenses.

        2. **Scenario Design**
           - Each scenario is defined by annual shocks to revenue and expenses.
           - Shocks are converted to quarterly growth rates and applied iteratively.

        3. **Capital & RWA**
           - RWA is proxied as a function of operating expenses.
           - Capital evolves by absorbing a share of losses after stress.
           - A floor is applied to enforce a minimum capital ratio.

        4. **Risk Appetite Metrics**
           - Capital ratio is compared against a user-defined threshold.
           - Efficiency ratio (expense / revenue) is assessed against a maximum tolerance.
           - Breaches are counted over the projection horizon.

        ### Controls & Limitations

        - This is a simplified educational model, not a production regulatory engine.
        - All figures are synthetic and for demonstration only.
        - Assumptions can be adjusted via the sidebar to perform what-if analysis.
        """
    )


# -------------------------------------------------------------------
# 6. Streamlit app
# -------------------------------------------------------------------

def main():
    st.title("Financial Stress Testing & Scenario Analysis Dashboard")
    st.caption(
        "Synthetic multi-scenario projections with capital and efficiency metrics, "
        "built for demonstration and interview discussion."
    )

    # Data
    df_hist = generate_historical_data()

    # Sidebar controls
    st.sidebar.header("Model Settings")

    horizon = st.sidebar.slider(
        "Projection horizon (quarters)",
        min_value=4,
        max_value=12,
        value=9,
        step=1,
    )

    cap_threshold = st.sidebar.slider(
        "Risk appetite: minimum capital ratio",
        min_value=0.08,
        max_value=0.16,
        value=0.11,
        step=0.005,
    )

    eff_threshold = st.sidebar.slider(
        "Risk appetite: max efficiency ratio (expense / revenue)",
        min_value=0.70,
        max_value=1.10,
        value=0.95,
        step=0.01,
    )

    capital_floor_ratio = st.sidebar.slider(
        "Regulatory capital floor ratio",
        min_value=0.06,
        max_value=0.12,
        value=0.08,
        step=0.005,
    )

    st.sidebar.markdown("---")
    st.sidebar.subheader("Scenario Assumptions (annual shocks)")

    with st.sidebar.expander("Base Case"):
        base_rev = st.number_input(
            "Base case revenue growth (annual)",
            min_value=-0.20,
            max_value=0.20,
            value=-0.01,
            step=0.01,
            format="%.2f",
            key="base_rev",
        )
        base_exp = st.number_input(
            "Base case expense growth (annual)",
            min_value=-0.10,
            max_value=0.20,
            value=0.01,
            step=0.01,
            format="%.2f",
            key="base_exp",
        )

    with st.sidebar.expander("Moderate Stress"):
        mod_rev = st.number_input(
            "Moderate stress revenue shock (annual)",
            min_value=-0.50,
            max_value=0.10,
            value=-0.03,
            step=0.01,
            format="%.2f",
            key="mod_rev",
        )
        mod_exp = st.number_input(
            "Moderate stress expense growth (annual)",
            min_value=-0.10,
            max_value=0.30,
            value=0.02,
            step=0.01,
            format="%.2f",
            key="mod_exp",
        )

    with st.sidebar.expander("Severe Stress"):
        sev_rev = st.number_input(
            "Severe stress revenue shock (annual)",
            min_value=-0.80,
            max_value=0.05,
            value=-0.08,
            step=0.01,
            format="%.2f",
            key="sev_rev",
        )
        sev_exp = st.number_input(
            "Severe stress expense growth (annual)",
            min_value=-0.10,
            max_value=0.40,
            value=0.03,
            step=0.01,
            format="%.2f",
            key="sev_exp",
        )

    scenario_params = {
        "Base Case": {"rev_shock_yr": base_rev, "exp_shock_yr": base_exp},
        "Moderate Stress": {"rev_shock_yr": mod_rev, "exp_shock_yr": mod_exp},
        "Severe Stress": {"rev_shock_yr": sev_rev, "exp_shock_yr": sev_exp},
    }

    # Run projections
    df_proj = build_all_scenarios(
        df_hist,
        scenario_params=scenario_params,
        horizon_quarters=horizon,
        capital_floor_ratio=capital_floor_ratio,
    )

    tab_dash, tab_doc = st.tabs(["Dashboard", "Model Documentation"])

    # ---------------- Dashboard tab ----------------
    with tab_dash:
        scenario_selected = st.selectbox(
            "Select scenario to view in detail",
            list(scenario_params.keys()),
        )

        agg_by_scenario = {}
        metrics_by_scenario = {}
        for name in scenario_params.keys():
            agg, metrics = summarize_scenario(
                df_proj, name, cap_threshold=cap_threshold, eff_threshold=eff_threshold
            )
            agg_by_scenario[name] = agg
            metrics_by_scenario[name] = metrics

        agg_sel = agg_by_scenario[scenario_selected]
        m_sel = metrics_by_scenario[scenario_selected]

        # KPI cards
        col1, col2, col3 = st.columns(3)
        col1.metric(
            "Cumulative stressed losses",
            f"${m_sel['cumulative_loss']:,.0f}",
        )
        col2.metric(
            "Minimum capital ratio",
            f"{m_sel['min_capital_ratio']:.2%}",
        )
        col3.metric(
            "# risk appetite breaches",
            int(m_sel["breaches"]),
        )

        # Explanation of breach definition
        st.caption(
            "📌 A 'risk appetite breach' is counted when either the quarter-end "
            "capital ratio falls below the minimum capital ratio threshold or the "
            "efficiency ratio (expense / revenue) rises above its maximum threshold "
            "set in the sidebar."
        )

        # Charts
        st.plotly_chart(
            plot_capital_ratio(agg_sel, cap_threshold, scenario_selected),
            use_container_width=True,
        )

        st.plotly_chart(
            plot_pretax_income_by_lob(df_proj, scenario_selected),
            use_container_width=True,
        )

        # Detailed table
        st.subheader("Projected summary – firm level")
        st.dataframe(
            agg_sel.style.format(
                {
                    "total_expense": "{:,.0f}",
                    "allocated_revenue": "{:,.0f}",
                    "pretax_income": "{:,.0f}",
                    "capital": "{:,.0f}",
                    "rwa": "{:,.0f}",
                    "capital_ratio": "{:.2%}",
                    "efficiency_ratio": "{:.2%}",
                }
            )
        )

        st.markdown("---")
        st.subheader("Scenario comparison – key metrics")

        comp_rows = []
        for name, m in metrics_by_scenario.items():
            comp_rows.append(
                {
                    "Scenario": name,
                    "Cumulative Loss": m["cumulative_loss"],
                    "Min Capital Ratio": m["min_capital_ratio"],
                    "# Breaches": m["breaches"],
                }
            )
        df_comp = pd.DataFrame(comp_rows)

        st.dataframe(
            df_comp.style.format(
                {
                    "Cumulative Loss": "{:,.0f}",
                    "Min Capital Ratio": "{:.2%}",
                }
            )
        )

        st.plotly_chart(
            plot_capital_ratio_multi(agg_by_scenario, cap_threshold),
            use_container_width=True,
        )

        # Alerts / escalation
        st.markdown("---")
        st.subheader("Alerts & escalation flags")

        any_breach = False
        for name, m in metrics_by_scenario.items():
            if m["breaches"] > 0 or m["min_capital_ratio"] < cap_threshold:
                any_breach = True
                st.error(
                    f"{name}: {m['breaches']} breach(es), "
                    f"minimum capital ratio {m['min_capital_ratio']:.2%} "
                    f"(threshold {cap_threshold:.2%})."
                )

        if not any_breach:
            st.success("No Risk Appetite breaches detected under any scenario.")

        # Downloads
        st.markdown("---")
        st.subheader("Download reports")

        col_a, col_b = st.columns(2)

        if PPTX_AVAILABLE:
            ppt_bytes = build_ppt_report(agg_by_scenario, metrics_by_scenario)
            col_a.download_button(
                "Download PowerPoint summary",
                data=ppt_bytes,
                file_name="stress_testing_summary.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        else:
            col_a.info("Install 'python-pptx' to enable PowerPoint export.")

        if PDF_AVAILABLE:
            pdf_bytes = build_pdf_report(
                metrics_by_scenario,
                cap_threshold=cap_threshold,
                eff_threshold=eff_threshold,
            )
        else:
            col_b.info("Install 'reportlab' to enable PDF export.")

        with st.expander("Show sample of raw historical data"):
            st.dataframe(df_hist.head(50))

    # ---------------- Documentation tab ----------------
    with tab_doc:
        make_documentation_tab()


if __name__ == "__main__":
    main()
